from pptx import Presentation
from app.utils.markdown_utils import to_markdown
from app.utils import get_logger


class PPTParser:
    def __init__(self):
        self.logger = get_logger(__name__)

    def _extract_slide_text(self, slide, slide_index: int) -> str:
        """Trích xuất toàn bộ text trong 1 slide."""
        texts = []
        for shape_index, shape in enumerate(slide.shapes, start=1):
            try:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                    self.logger.debug(f"📝 Slide {slide_index} - Shape {shape_index}: {shape.text.strip()[:60]}...")
            except Exception as e:
                self.logger.warning(f"⚠️ Lỗi khi đọc shape {shape_index} ở slide {slide_index}: {e}")
        return "\n".join(texts)


    def parse(self, file_path: str) -> str:
        """Phân tích file PowerPoint (PPTX) và trích xuất toàn bộ nội dung dạng Markdown."""
        self.logger.info(f"📊 Bắt đầu parsing PPTX: {file_path}")

        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        self.logger.info(f"🔍 Tệp có {total_slides} slide.")

        slides_content = []
        for i, slide in enumerate(prs.slides, start=1):
            self.logger.debug(f"➡️ Đang xử lý slide {i}/{total_slides}")
            slide_text = self._extract_slide_text(slide, i)
            if not slide_text.strip():
                self.logger.debug(f"⚪ Slide {i} trống hoặc không chứa text.")
            slides_content.append(f"## Slide {i}\n{slide_text}")

        full_text = "\n\n".join(slides_content)
        md = to_markdown(full_text)

        self.logger.info(f"✅ Hoàn tất parsing PPTX: {file_path}")
        return md
